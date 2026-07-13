#!/usr/bin/env python3
import argparse
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


def clean_text(value: str) -> str:
    return " ".join(value.replace("\xa0", " ").split())


def escape_table_cell(value: str) -> str:
    return clean_text(value).replace("|", "\\|")


def paragraph_to_markdown(paragraph) -> str:
    text = clean_text(paragraph.text)
    if not text:
        return ""

    style_name = (getattr(paragraph.style, "name", "") or "").lower()
    if "heading 1" in style_name or "标题 1" in style_name:
        return f"# {text}"
    if "heading 2" in style_name or "标题 2" in style_name:
        return f"## {text}"
    if "heading 3" in style_name or "标题 3" in style_name:
        return f"### {text}"
    if "heading 4" in style_name or "标题 4" in style_name:
        return f"#### {text}"
    if "heading 5" in style_name or "标题 5" in style_name:
        return f"##### {text}"
    if "heading 6" in style_name or "标题 6" in style_name:
        return f"###### {text}"
    if "list bullet" in style_name or "项目符号" in style_name:
        return f"- {text}"
    if "list number" in style_name or "编号" in style_name:
        return f"1. {text}"
    return text


def table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [escape_table_cell(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""

    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def parse_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("缺少 python-docx 依赖，请先安装文档解析依赖") from exc

    document = Document(path)
    parts = []

    for paragraph in document.paragraphs:
        markdown = paragraph_to_markdown(paragraph)
        if markdown:
            parts.append(markdown)

    for table in document.tables:
        markdown = table_to_markdown(table)
        if markdown:
            parts.append(markdown)

    return "\n\n".join(parts)


def normalize_pdf_text(text: str) -> str:
    lines = []
    blank = False
    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = raw_line.rstrip()
        stripped = clean_text(line)
        if not stripped:
            if not blank and lines:
                lines.append("")
                blank = True
            continue
        blank = False
        if stripped.startswith(("• ", "● ", "· ")):
            stripped = "- " + stripped[2:].strip()
        lines.append(stripped)
    return "\n".join(lines).strip()


def parse_pdf(path: str) -> str:
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("缺少 pdfplumber 依赖，请先安装文档解析依赖") from exc

    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text(layout=True) or page.extract_text() or ""
            text = normalize_pdf_text(text)
            if text:
                parts.append(text)

    return "\n\n".join(parts)


def parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx 依赖，请先安装文档解析依赖") from exc

    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = [f"# 第 {index} 页"]
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                markdown = table_to_markdown(shape.table)
                if markdown:
                    parts.append(markdown)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = clean_text(paragraph.text)
                if not text:
                    continue
                if getattr(paragraph, "level", 0) > 0 or paragraph.text.startswith(("•", "-")):
                    parts.append(f"- {text.lstrip('•- ')}")
                else:
                    parts.append(text)
        if len(parts) > 1:
            slides.append("\n\n".join(parts))
    return "\n\n".join(slides)


def main() -> int:
    parser = argparse.ArgumentParser(description="提取 docx/pdf/pptx 文档正文")
    parser.add_argument("--type", required=True, choices=["docx", "pdf", "pptx"])
    parser.add_argument("--file", required=True)
    args = parser.parse_args()

    try:
        if args.type == "docx":
            content = parse_docx(args.file)
        elif args.type == "pdf":
            content = parse_pdf(args.file)
        else:
            content = parse_pptx(args.file)
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        return 1

    sys.stdout.write(content)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
